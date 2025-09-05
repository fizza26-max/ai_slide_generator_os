
import io, os, json
from typing import List, Dict
import streamlit as st

from transformers import pipeline
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ---------------- Helpers -----------------

@st.cache_resource(show_spinner=False)
def load_model(model_name: str = "google/flan-t5-small"):
    # single pipeline for text2text generation to handle bullets + speaker notes
    return pipeline("text2text-generation", model=model_name, tokenizer=model_name, device_map=None)

def chunk_text(text: str, max_chars: int = 900) -> List[str]:
    text = " ".join(text.split())
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + max_chars, len(text))
        slice_ = text[start:end]
        last_punct = max(slice_.rfind("."), slice_.rfind("?"), slice_.rfind("!"))
        if last_punct == -1 or end == len(text):
            cut = end
        else:
            cut = start + last_punct + 1
        chunks.append(text[start:cut].strip())
        start = cut
    merged=[]
    for ch in chunks:
        if merged and len(ch) < 200:
            merged[-1] = merged[-1] + " " + ch
        else:
            merged.append(ch)
    return [c for c in merged if c]

def generate_bullets(model, text: str, max_bullets: int = 5) -> List[str]:
    prompt = ("Summarize the following content into concise presentation bullets. "
              f"Return {max_bullets} or fewer bullet points. Do not number them. Keep each bullet under 18 words.\\n\\nContent:\\n" + text)
    out = model(prompt, max_new_tokens=160, do_sample=False)[0]["generated_text"]
    lines = [ln.strip("-• \n\t") for ln in out.split("\n") if ln.strip()]
    if len(lines) <= 2:
        tmp=[]
        for piece in out.replace("•","\\n").split("\\n"):
            for seg in piece.split("."):
                seg = seg.strip(" -•\\t")
                if len(seg.split()) >= 3:
                    tmp.append(seg)
        lines = tmp
    bullets = []
    for ln in lines:
        ln = ln.lstrip("0123456789). ").strip(" -•\\t")
        if ln:
            bullets.append(ln)
        if len(bullets) >= max_bullets:
            break
    return bullets

def generate_speaker_notes(model, title: str, bullets: List[str]) -> str:
    """
    Create a short speaker note for a slide based on title + bullets.
    """
    text = f"Write a concise speaker note for a presentation slide titled '{title}'. Use the following bullet points as the slide content: " + " | ".join(bullets)
    prompt = text + " Keep it under 120 words and conversational."
    out = model(prompt, max_new_tokens=160, do_sample=False)[0]["generated_text"]
    return out.strip()

def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2],16), int(hex_color[2:4],16), int(hex_color[4:6],16))

def pick_contrasting_text(rgb: RGBColor) -> RGBColor:
    brightness = (0.299*rgb[0] + 0.587*rgb[1] + 0.114*rgb[2])
    return RGBColor(0,0,0) if brightness > 186 else RGBColor(255,255,255)

def add_title_slide(prs: Presentation, title: str, subtitle: str, theme_hex: str, font_name: str):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme_hex)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    title_shape.text = title
    subtitle_shape.text = subtitle
    tcolor = pick_contrasting_text(hex_to_rgb(theme_hex))
    for shp in (title_shape, subtitle_shape):
        for p in shp.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = tcolor
                try:
                    run.font.name = font_name
                except Exception:
                    pass

def add_bullet_slide(prs: Presentation, title: str, bullets: List[str], speaker_note: str, theme_hex: str, font_name: str):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme_hex)
    title_shape = slide.shapes.title
    title_shape.text = title
    tcolor = pick_contrasting_text(hex_to_rgb(theme_hex))
    for p in title_shape.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = tcolor
            run.font.bold = True
            try:
                run.font.name = font_name
            except Exception:
                pass
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i,b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(22)
            run.font.color.rgb = tcolor
            try:
                run.font.name = font_name
            except Exception:
                pass
    # add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.clear()
    text_frame.text = speaker_note

def make_presentation(slides_data: List[Dict], theme_hex: str = "#1F2937", font_name: str = "Calibri") -> bytes:
    prs = Presentation()
    for s in slides_data:
        if s["type"] == "title":
            add_title_slide(prs, s.get("title","Title"), s.get("subtitle",""), theme_hex, font_name)
        else:
            add_bullet_slide(prs, s.get("title","Slide"), s.get("bullets",[]), s.get("notes",""), theme_hex, font_name)
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()

def build_google_slides_requests(slides_data: List[Dict], theme_hex: str, font_family: str) -> Dict:
    """
    Build a Google Slides batchUpdate-style requests list that can be sent to the Google Slides API.
    This returns a JSON dict with 'requests' that the user can apply after obtaining OAuth credentials.
    """
    requests = []
    # create presentation (title)
    # For simplicity: we'll create slides and add text boxes. Users will need to run this via Slides API.
    for idx, s in enumerate(slides_data):
        # create slide
        create_req = {
            "createSlide": {
                "slideLayoutReference": {
                    "predefinedLayout": "TITLE_AND_BODY"
                },
                "insertionIndex": idx
            }
        }
        requests.append(create_req)
        # add title text
        title_text = s.get("title","")
        body_text = ""
        if s["type"] == "title":
            body_text = s.get("subtitle","")
        else:
            body_text = "\\n".join(s.get("bullets",[]))
        # text insertion requests (for newly created slide) using placeholders
        requests.append({
            "insertText": {
                "objectId": f"TITLE_{idx}",
                "insertionIndex": 0,
                "text": title_text
            }
        })
        requests.append({
            "insertText": {
                "objectId": f"BODY_{idx}",
                "insertionIndex": 0,
                "text": body_text
            }
        })
    return {"requests": requests, "note": "This JSON contains placeholder requests. You must map objectIds after creation via the Slides API when running batchUpdate."}

# ---------------- Streamlit UI -----------------

st.set_page_config(page_title="AI Slide Generator — Themes + Speaker Notes + Google Slides", layout="wide")

st.title("AI Slide Generator — Themes, Speaker Notes & Google Slides Export")
st.write("Open-source models only. This app produces `.pptx` files with speaker notes and a Google Slides batchUpdate JSON you can use with the Slides API (needs OAuth).")

with st.sidebar:
    st.header("Settings")
    model_choice = st.selectbox("Model", ["google/flan-t5-small","google/flan-t5-base"], index=0)
    max_bullets = st.slider("Bullets per slide", 3, 6, 4)
    max_chars = st.slider("Max chars per chunk", 600, 1400, 900, step=100)
    add_conclusion = st.checkbox("Add conclusion slide", True)
    show_preview = st.checkbox("Show preview", True)
    st.markdown("### Themes (preset palettes & fonts)")
    theme_preset = st.selectbox("Choose theme preset", ["Classic Dark","Light Minimal","Blue Gradient","Warm Accent"])
    font_choice = st.selectbox("Font family (used for PPTX)", ["Calibri","Arial","Georgia","Tahoma","Verdana"])

col1, col2 = st.columns([2,1])
with col1:
    text_input = st.text_area("Paste source text", height=320)
with col2:
    title_override = st.text_input("Presentation Title", "")
    subtitle_override = st.text_input("Subtitle (author/date)", "")
    st.markdown("**Quick actions**")
    sample = st.button("Load sample text")

if sample:
    text_input = ("Large Language Models (LLMs) are transforming the way people interact with technology. "
    "They enable natural language interfaces that understand context, follow instructions, and generate useful content. "
    "This demo shows automatic chunking, bullet generation, speaker notes, and export to PPTX or Google Slides JSON.")

gen = st.button("Generate Slides")

theme_map = {
    "Classic Dark": "#1F2937",
    "Light Minimal": "#FFFFFF",
    "Blue Gradient": "#0B69FF",
    "Warm Accent": "#FF7043"
}

if gen:
    if not text_input.strip():
        st.error("Please paste text first")
        st.stop()
    with st.spinner("Loading model..."):
        model = load_model(model_choice)
    chunks = chunk_text(text_input, max_chars=max_chars)
    slides = []
    title = title_override.strip() or (text_input.splitlines()[0] if text_input.splitlines() and len(text_input.splitlines()[0])>6 else "Presentation")
    subtitle = subtitle_override.strip() or "Generated automatically"
    slides.append({"type":"title","title":title,"subtitle":subtitle})
    for i,ch in enumerate(chunks, start=1):
        bullets = generate_bullets(model, ch, max_bullets=max_bullets)
        slide_title = f"Section {i}"
        notes = generate_speaker_notes(model, slide_title, bullets)
        slides.append({"type":"content","title":slide_title,"bullets":bullets,"notes":notes})
    if add_conclusion:
        bullets = generate_bullets(model, "Overall: " + text_input, max_bullets=4)
        notes = generate_speaker_notes(model, "Conclusion", bullets)
        slides.append({"type":"content","title":"Conclusion","bullets":bullets,"notes":notes})

    # Preview
    if show_preview:
        st.subheader("Preview")
        for s in slides:
            if s["type"]=="title":
                st.markdown(f"### {s['title']}"); st.markdown(f"*{s['subtitle']}*")
            else:
                st.markdown(f"**{s['title']}**")
                for b in s["bullets"]:
                    st.markdown(f"- {b}")
                st.markdown(f"_Speaker note:_ {s.get('notes','')}")

    # Build PPTX
    theme_hex = theme_map.get(theme_preset, "#1F2937")
    with st.spinner("Building PPTX..."):
        pptx_bytes = make_presentation(slides, theme_hex=theme_hex, font_name=font_choice)
    st.success("PPTX ready")
    st.download_button("Download .pptx", data=pptx_bytes, file_name="ai_slides_with_notes.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    # Build Google Slides JSON (requests)
    gs_json = build_google_slides_requests(slides, theme_hex, font_choice)
    gs_bytes = io.BytesIO(json.dumps(gs_json, indent=2).encode("utf-8"))
    st.download_button("Download Google Slides batchUpdate JSON (manual apply)", data=gs_bytes, file_name="google_slides_requests.json", mime="application/json")
    st.info("Note: The JSON contains placeholder requests. To fully apply them you must run the Google Slides API batchUpdate and map objectIds created at runtime. See instructions below.")

    st.markdown("### How to apply the Google Slides JSON (quick guide)")
    st.markdown("""
1. Create a Google Cloud project and enable the Google Slides API.
2. Obtain OAuth credentials (service account or OAuth client). For service accounts, share the target Google Drive folder/presentation with the service account email.
3. Use the Slides API to create a presentation, capture the slide/object ids returned, then adjust the `insertText` requests objectIds to the real shapes before calling `batchUpdate`.
4. If you prefer, use the Google API Python client to automate these steps. This app intentionally outputs requests in a simplified format to be safe for users to adapt.
    """)
else:
    st.info("Paste text and click Generate Slides.")
